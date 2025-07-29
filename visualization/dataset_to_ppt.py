import os
from pptx import Presentation
from pptx.util import Inches

def create_ppt_from_dirs(image_dirs, ppt_path):
    """
    Create a PowerPoint file where each slide contains one comparison image from the given directories.
    :param image_dirs: List of directories containing comparison JPGs
    :param ppt_path: Output path for the .pptx file
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for img_dir in image_dirs:
        # Sort images for consistent order
        for img_file in sorted(os.listdir(img_dir)):
            if not img_file.lower().endswith('.jpg'):
                continue
            img_path = os.path.join(img_dir, img_file)
            slide = prs.slides.add_slide(blank_layout)
            # Add picture with margins
            slide.shapes.add_picture(
                img_path,
                Inches(0.5),  # left margin
                Inches(0.5),  # top margin
                width=Inches(9)  # fit width
            )
    prs.save(ppt_path)
    print(f"PowerPoint saved to {ppt_path}")

if __name__ == '__main__':
    # Directories where comparison images are already generated
    p0_comparisons = '/home/ubuntu/Desktop/JY/ddrm/ultrasound/datasets_v0.03/P0/comparisons'
    p2_comparisons = '/home/ubuntu/Desktop/JY/ddrm/ultrasound/datasets_v0.03/P2/comparisons'
    # Output PPT path
    output_ppt = '/home/ubuntu/Desktop/JY/ddrm/comparisons_slides.pptx'

    create_ppt_from_dirs([p0_comparisons, p2_comparisons], output_ppt)
